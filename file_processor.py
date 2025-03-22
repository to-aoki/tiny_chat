import io
import re
from urllib.parse import urlparse
from bs4 import BeautifulSoup


class PDFProcessor:
    """PDFファイル処理クラス"""

    @staticmethod
    def extract_text_from_bytes(pdf_bytes):
        """
        PDFバイトデータからテキストを抽出する基本メソッド

        Args:
            pdf_bytes: PDFファイルのバイトデータ

        Returns:
            tuple: (抽出されたテキスト, ページ数, エラーメッセージ)
        """
        try:
            from pypdf import PdfReader

            pdf_stream = io.BytesIO(pdf_bytes)
            pdf_reader = PdfReader(pdf_stream)

            extracted_text = ""
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                extracted_text += f"\n--- ページ {page_num + 1} ---\n{page_text}"

            return extracted_text, len(pdf_reader.pages), None
        except Exception as e:
            return None, 0, str(e)

    @staticmethod
    def extract_pdf_text(pdf_file):
        """
        PDFファイルからテキストを抽出 (Streamlit用)

        Args:
            pdf_file: PDFファイルオブジェクト（Streamlitのアップローダーから取得）

        Returns:
            tuple: (抽出されたテキスト, ページ数, エラーメッセージ)
        """
        try:
            pdf_bytes = pdf_file.getvalue()
            return PDFProcessor.extract_text_from_bytes(pdf_bytes)
        except Exception as e:
            return None, 0, str(e)


class ExcelProcessor:
    """Excelファイル処理クラス"""

    @staticmethod
    def extract_text_from_bytes(excel_bytes, include_sheet_names=True):
        """
        Excelバイトデータからテキストを抽出する基本メソッド

        Args:
            excel_bytes: Excelファイルのバイトデータ
            include_sheet_names (bool): シート名を含めるかどうか

        Returns:
            tuple: (抽出されたテキスト, シート数, エラーメッセージ)
        """
        try:
            import openpyxl

            excel_stream = io.BytesIO(excel_bytes)

            # openpyxlでExcelファイルを読み込む
            workbook = openpyxl.load_workbook(excel_stream, data_only=True)
            sheet_names = workbook.sheetnames

            extracted_text = ""
            for sheet_name in sheet_names:
                if include_sheet_names:
                    extracted_text += f"\n--- シート: {sheet_name} ---\n"

                sheet = workbook[sheet_name]

                # 各行をテキストとして処理
                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        cell_value = cell.value
                        if cell_value is not None:
                            row_values.append(str(cell_value))

                    # 行の値を1つの文として結合
                    if row_values:
                        row_text = " ".join(row_values)
                        extracted_text += row_text + "\n"

                extracted_text += "\n"

            return extracted_text, len(sheet_names), None
        except Exception as e:
            return None, 0, str(e)

    @staticmethod
    def extract_excel_text(excel_file, include_sheet_names=True):
        """
        Excelファイルからテキストを抽出 (Streamlit用)

        Args:
            excel_file: Excelファイルオブジェクト（Streamlitのアップローダーから取得）
            include_sheet_names (bool): シート名を含めるかどうか

        Returns:
            tuple: (抽出されたテキスト, シート数, エラーメッセージ)
        """
        try:
            excel_bytes = excel_file.getvalue()
            return ExcelProcessor.extract_text_from_bytes(excel_bytes, include_sheet_names)
        except Exception as e:
            return None, 0, str(e)


class WordProcessor:
    """Wordファイル処理クラス"""

    @staticmethod
    def extract_text_from_bytes(word_bytes):
        """
        Wordバイトデータからテキストを抽出する基本メソッド

        Args:
            word_bytes: Wordファイルのバイトデータ

        Returns:
            tuple: (抽出されたテキスト, エラーメッセージ)
        """
        try:
            import docx

            word_stream = io.BytesIO(word_bytes)
            doc = docx.Document(word_stream)

            extracted_text = ""

            # パラグラフからテキストを抽出
            for para in doc.paragraphs:
                if para.text.strip():  # 空のパラグラフをスキップ
                    extracted_text += para.text + "\n"

            # テーブルからテキストを抽出
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    extracted_text += row_text + "\n"
                extracted_text += "\n"

            return extracted_text, None
        except Exception as e:
            return None, str(e)

    @staticmethod
    def extract_word_text(word_file):
        """
        Wordファイルからテキストを抽出 (Streamlit用)

        Args:
            word_file: Wordファイルオブジェクト（Streamlitのアップローダーから取得）

        Returns:
            tuple: (抽出されたテキスト, エラーメッセージ)
        """
        try:
            word_bytes = word_file.getvalue()
            return WordProcessor.extract_text_from_bytes(word_bytes)
        except Exception as e:
            return None, str(e)


class PowerPointProcessor:
    """PowerPointファイル処理クラス"""

    @staticmethod
    def extract_text_from_bytes(pptx_bytes):
        """
        PowerPointバイトデータからテキストを抽出する基本メソッド

        Args:
            pptx_bytes: PowerPointファイルのバイトデータ

        Returns:
            tuple: (抽出されたテキスト, スライド数, エラーメッセージ)
        """
        try:
            from pptx import Presentation

            pptx_stream = io.BytesIO(pptx_bytes)
            prs = Presentation(pptx_stream)

            extracted_text = ""
            slide_count = len(prs.slides)

            for i, slide in enumerate(prs.slides):
                extracted_text += f"\n--- スライド {i + 1} ---\n"

                # スライドのタイトル
                if slide.shapes.title:
                    title_text = slide.shapes.title.text
                    extracted_text += f"タイトル: {title_text}\n"

                # スライド内のテキスト
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        extracted_text += shape.text + "\n"

                extracted_text += "\n"

            return extracted_text, slide_count, None
        except Exception as e:
            return None, 0, str(e)

    @staticmethod
    def extract_pptx_text(pptx_file):
        """
        PowerPointファイルからテキストを抽出 (Streamlit用)

        Args:
            pptx_file: PowerPointファイルオブジェクト（Streamlitのアップローダーから取得）

        Returns:
            tuple: (抽出されたテキスト, スライド数, エラーメッセージ)
        """
        try:
            pptx_bytes = pptx_file.getvalue()
            return PowerPointProcessor.extract_text_from_bytes(pptx_bytes)
        except Exception as e:
            return None, 0, str(e)


class TextFileProcessor:
    """テキストファイル処理クラス"""

    @staticmethod
    def extract_text_from_bytes(text_bytes):
        """
        テキストバイトデータからテキストを抽出する基本メソッド

        Args:
            text_bytes: テキストファイルのバイトデータ

        Returns:
            tuple: (抽出されたテキスト, エラーメッセージ)
        """
        try:
            # テキストファイルの場合、エンコーディングの検出を試みる
            # まずUTF-8で試す
            try:
                extracted_text = text_bytes.decode('utf-8')
                return extracted_text, None
            except UnicodeDecodeError:
                pass

            try:
                extracted_text = text_bytes.decode('shift_jis')
                return extracted_text, None
            except UnicodeDecodeError:
                pass

            try:
                extracted_text = text_bytes.decode('cp932')
                return extracted_text, None
            except UnicodeDecodeError:
                return None, "テキストファイルのエンコーディングを認識できませんでした"

        except Exception as e:
            return None, str(e)

    @staticmethod
    def extract_text(text_file):
        """
        テキストファイルからテキストを抽出 (Streamlit用)

        Args:
            text_file: テキストファイルオブジェクト（Streamlitのアップローダーから取得）

        Returns:
            tuple: (抽出されたテキスト, エラーメッセージ)
        """
        try:
            text_bytes = text_file.getvalue()
            return TextFileProcessor.extract_text_from_bytes(text_bytes)
        except Exception as e:
            return None, str(e)


class HTMLProcessor:
    """HTMLコンテンツ処理クラス"""

    @staticmethod
    def extract_text_from_bytes(html_bytes):
        """
        HTMLバイトデータからテキストを抽出する基本メソッド

        Args:
            html_bytes: HTMLコンテンツのバイトデータ

        Returns:
            tuple: (抽出されたテキスト, メッセージ)
        """
        try:
            # エンコーディングの検出を試みる
            try:
                html_content = html_bytes.decode('utf-8')
            except UnicodeDecodeError:
                try:
                    html_content = html_bytes.decode('shift_jis')
                except UnicodeDecodeError:
                    try:
                        html_content = html_bytes.decode('cp932')
                    except UnicodeDecodeError:
                        return None, "HTMLのエンコーディングを認識できませんでした"

            soup = BeautifulSoup(html_content, 'html.parser')

            # スクリプトと不要なタグを削除
            for script in soup(["script", "style", "footer", "nav"]):
                script.extract()

            # テキストを抽出して整形
            text = soup.get_text(separator='\n')
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            extracted_text = '\n'.join(chunk for chunk in chunks if chunk)

            # タイトルを取得
            title = soup.title.string if soup.title else "Webページ"

            return extracted_text, f"HTMLページ '{title}' からテキストを抽出しました"
        except Exception as e:
            return None, str(e)


class FileProcessorFactory:
    """ファイル処理ファクトリークラス"""

    @staticmethod
    def get_processor(file_extension):
        """
        ファイル拡張子に基づいて適切なプロセッサを返す

        Args:
            file_extension (str): ファイル拡張子（.pdf, .xlsx, .docx, .pptx, .txt など）

        Returns:
            適切なプロセッサクラス
        """
        extension_map = {
            '.pdf': PDFProcessor,
            '.xlsx': ExcelProcessor,
            '.xls': ExcelProcessor,
            '.docx': WordProcessor,
            '.pptx': PowerPointProcessor,
            '.txt': TextFileProcessor,
            '.csv': TextFileProcessor,
            '.json': TextFileProcessor,
            '.md': TextFileProcessor,
            '.html': HTMLProcessor,
            '.htm': HTMLProcessor,
        }

        return extension_map.get(file_extension.lower())


class URIProcessor:
    """URI処理クラス"""

    @staticmethod
    def detect_uri(text):
        """
        テキスト内のURIを検出

        Args:
            text (str): 検査するテキスト

        Returns:
            list: 検出されたURIのリスト
        """
        uri_pattern = r'https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+'
        return re.findall(uri_pattern, text)

    @staticmethod
    def process_uri(uri, max_length=4000):
        """
        URIからコンテンツを取得して処理

        Args:
            uri (str): 処理するURI
            max_length (int): 最大コンテキスト長

        Returns:
            tuple: (抽出されたコンテンツ, メッセージ)
        """
        try:
            # URLの判定
            parsed_url = urlparse(uri)
            if not parsed_url.scheme or not parsed_url.netloc:
                return None, f"無効なURI形式です: {uri}"

            # リクエスト実行
            import requests
            response = requests.get(uri, stream=True)
            if response.status_code != 200:
                return None, f"URIからのコンテンツ取得に失敗しました（ステータスコード: {response.status_code}）"

            # Content-Typeを確認
            content_type = response.headers.get('Content-Type', '').lower()
            content = response.content

            # 各コンテンツタイプごとに処理
            if 'application/pdf' in content_type:
                extract_text, page_count, error = PDFProcessor.extract_text_from_bytes(content)
                message = f"PDFから{page_count}ページのテキストを抽出しました" if not error else error

            elif 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' in content_type or 'application/vnd.ms-excel' in content_type:
                extract_text, sheet_count, error = ExcelProcessor.extract_text_from_bytes(content)
                message = f"Excelから{sheet_count}シートのテキストを抽出しました" if not error else error

            elif 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' in content_type or 'application/msword' in content_type:
                extract_text, error = WordProcessor.extract_text_from_bytes(content)
                message = "Wordドキュメントからテキストを抽出しました" if not error else error

            elif 'application/vnd.openxmlformats-officedocument.presentationml.presentation' in content_type or 'application/vnd.ms-powerpoint' in content_type:
                extract_text, slide_count, error = PowerPointProcessor.extract_text_from_bytes(content)
                message = f"PowerPointから{slide_count}スライドのテキストを抽出しました" if not error else error

            elif 'text/html' in content_type:
                extract_text, message = HTMLProcessor.extract_text_from_bytes(content)

            # その他のテキスト形式
            elif 'text/' in content_type:
                extract_text, error = TextFileProcessor.extract_text_from_bytes(content)
                message = "テキストコンテンツを取得しました" if not error else error

            else:
                extract_text, message = None, f"サポートされていないコンテンツタイプです: {content_type}"

            if extract_text:
                extract_text = extract_text[:max_length]

            return extract_text, message

        except Exception as e:
            return None, f"URLコンテンツ取得中にエラーが発生しました: {str(e)}"
